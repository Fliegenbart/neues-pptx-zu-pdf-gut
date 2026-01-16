"""
PPTX Parser
===========
Extrahiert ein semantisches SlideModel aus PowerPoint-Dateien.

Kernaufgaben:
1. Shapes parsen und klassifizieren
2. Lesereihenfolge heuristisch bestimmen
3. Bilder extrahieren (mit Hash für Caching)
4. Tabellen erkennen (echte + "Fake-Tables")
"""

import hashlib
import re
from pathlib import Path
from typing import Optional
from zipfile import ZipFile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table as PptxTable

from .models import (
    SlideModel, Slide, Block, BlockType, 
    Paragraph, TextRun, Table, TableCell,
    Figure, BoundingBox, ListStyle
)


class PPTXParser:
    """
    Parst PPTX-Dateien zu einem semantischen SlideModel.
    
    Usage:
        parser = PPTXParser()
        model = parser.parse("presentation.pptx")
    """
    
    # Schwellwerte für Heuristiken
    HEADING_MIN_FONT_SIZE = 18  # pt
    TITLE_PLACEHOLDER_TYPES = {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.SUBTITLE,
    }
    
    def __init__(self, extract_images: bool = True):
        """
        Args:
            extract_images: Ob Bilder extrahiert werden sollen
        """
        self.extract_images = extract_images
        self._image_cache: dict[str, bytes] = {}
    
    def parse(self, pptx_path: Path | str) -> SlideModel:
        """
        Parst eine PPTX-Datei.
        
        Args:
            pptx_path: Pfad zur PPTX-Datei
            
        Returns:
            SlideModel mit allen Folien und Inhalten
        """
        pptx_path = Path(pptx_path)
        
        # Bilder vorab extrahieren (für Hashing)
        if self.extract_images:
            self._extract_media(pptx_path)
        
        # PPTX laden
        prs = Presentation(str(pptx_path))
        
        # Modell aufbauen
        model = SlideModel(
            source_file=pptx_path,
            language="de",  # TODO: Aus PPTX extrahieren
        )
        
        # Dokument-Metadaten
        if prs.core_properties:
            model.title = prs.core_properties.title
            model.author = prs.core_properties.author
            model.subject = prs.core_properties.subject
        
        # Folien parsen
        for slide_num, pptx_slide in enumerate(prs.slides, 1):
            slide = self._parse_slide(pptx_slide, slide_num)
            model.slides.append(slide)
        
        return model
    
    def _extract_media(self, pptx_path: Path):
        """Extrahiert alle Medien und berechnet Hashes."""
        self._image_cache.clear()
        
        with ZipFile(pptx_path, 'r') as zf:
            for name in zf.namelist():
                if name.startswith('ppt/media/'):
                    data = zf.read(name)
                    # Key ist Dateiname ohne Pfad
                    key = Path(name).name
                    self._image_cache[key] = data
    
    def _parse_slide(self, pptx_slide, slide_num: int) -> Slide:
        """Parst eine einzelne Folie."""
        slide = Slide(
            number=slide_num,
            width_mm=254.0,   # Standard 16:9
            height_mm=142.9,
        )
        
        # Alle Shapes sammeln
        shapes_with_order = []
        
        for shape in pptx_slide.shapes:
            block = self._parse_shape(shape)
            if block and not block.is_empty:
                # Position für Lesereihenfolge
                bbox = self._get_bounding_box(shape)
                block.bbox = bbox
                shapes_with_order.append((shape, block, bbox))
        
        # Lesereihenfolge bestimmen
        ordered_blocks = self._determine_reading_order(shapes_with_order)
        slide.blocks = ordered_blocks
        
        # Speaker Notes
        if pptx_slide.has_notes_slide:
            notes_frame = pptx_slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                slide.notes = notes_frame.text.strip()
        
        return slide
    
    def _parse_shape(self, shape: BaseShape) -> Optional[Block]:
        """
        Parst ein Shape zu einem semantischen Block.
        
        Erkennt automatisch den Typ basierend auf:
        - Shape-Typ (Bild, Tabelle, Text)
        - Placeholder-Typ (Titel, Untertitel, Content)
        - Formatierung (Schriftgröße, Bold, etc.)
        """
        # Bilder
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._parse_picture(shape)
        
        # Tabellen
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return self._parse_table(shape)
        
        # Embedded Charts/Diagramme
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return self._parse_chart(shape)
        
        # Text-Shapes
        if shape.has_text_frame:
            return self._parse_text_shape(shape)
        
        # Gruppen rekursiv verarbeiten
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # TODO: Gruppierte Shapes verarbeiten
            pass
        
        return None
    
    def _parse_text_shape(self, shape: BaseShape) -> Optional[Block]:
        """Parst ein Text-Shape und bestimmt den semantischen Typ."""
        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return None
        
        # Paragraphen extrahieren
        paragraphs = []
        max_font_size = 0
        has_bullets = False
        
        for pptx_para in text_frame.paragraphs:
            para = self._parse_paragraph(pptx_para)
            if para and not para.is_empty:
                paragraphs.append(para)
                
                # Font-Size tracken für Heading-Erkennung
                for run in para.runs:
                    if run.font_size and run.font_size > max_font_size:
                        max_font_size = run.font_size
                
                # Bullet-Erkennung
                if pptx_para.level > 0 or self._has_bullet(pptx_para):
                    has_bullets = True
        
        if not paragraphs:
            return None
        
        # Block-Typ bestimmen
        block_type, heading_level, list_style = self._classify_text_block(
            shape, paragraphs, max_font_size, has_bullets
        )
        
        return Block(
            block_type=block_type,
            reading_order=0,  # Wird später gesetzt
            paragraphs=paragraphs,
            heading_level=heading_level,
            list_style=list_style,
            source_shape_id=str(shape.shape_id),
        )
    
    def _parse_paragraph(self, pptx_para) -> Optional[Paragraph]:
        """Parst einen Absatz mit allen Runs."""
        runs = []
        
        for pptx_run in pptx_para.runs:
            if not pptx_run.text:
                continue
                
            run = TextRun(
                text=pptx_run.text,
                bold=pptx_run.font.bold or False,
                italic=pptx_run.font.italic or False,
                underline=pptx_run.font.underline or False,
                font_size=self._emu_to_pt(pptx_run.font.size) if pptx_run.font.size else None,
                font_name=pptx_run.font.name,
            )
            
            # Hyperlink
            if pptx_run.hyperlink and pptx_run.hyperlink.address:
                run.hyperlink = pptx_run.hyperlink.address
            
            runs.append(run)
        
        if not runs:
            return None
        
        # Alignment
        alignment_map = {
            1: "left",
            2: "center",
            3: "right",
            4: "justify",
        }
        alignment = alignment_map.get(pptx_para.alignment, "left") if pptx_para.alignment else "left"
        
        return Paragraph(
            runs=runs,
            alignment=alignment,
            level=pptx_para.level or 0,
        )
    
    def _classify_text_block(
        self, 
        shape: BaseShape, 
        paragraphs: list[Paragraph],
        max_font_size: float,
        has_bullets: bool
    ) -> tuple[BlockType, int, ListStyle]:
        """
        Klassifiziert einen Text-Block.
        
        Returns:
            Tuple von (BlockType, heading_level, list_style)
        """
        # 1. Placeholder-Typ prüfen (höchste Priorität)
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            
            if ph_type in self.TITLE_PLACEHOLDER_TYPES:
                return BlockType.HEADING, 1, ListStyle.NONE
            
            if ph_type == PP_PLACEHOLDER.SUBTITLE:
                return BlockType.HEADING, 2, ListStyle.NONE
        
        # 2. Listen erkennen
        if has_bullets:
            return BlockType.LIST, 0, ListStyle.BULLET
        
        # 3. Heading durch Formatierung erkennen
        if max_font_size >= self.HEADING_MIN_FONT_SIZE:
            # Heading-Level basierend auf Font-Size
            if max_font_size >= 32:
                level = 1
            elif max_font_size >= 24:
                level = 2
            elif max_font_size >= 20:
                level = 3
            else:
                level = 4
            
            return BlockType.HEADING, level, ListStyle.NONE
        
        # 4. Default: Paragraph
        return BlockType.PARAGRAPH, 0, ListStyle.NONE
    
    def _parse_picture(self, shape: Picture) -> Optional[Block]:
        """Parst ein Bild-Shape."""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Hash für Caching
            image_hash = hashlib.md5(image_bytes).hexdigest()
            
            # MIME-Type
            mime_type = f"image/{image.ext}"
            if image.ext == "jpg":
                mime_type = "image/jpeg"
            
            # Alt-Text aus PPTX (falls vorhanden)
            existing_alt = None
            if hasattr(shape, '_element'):
                # Versuche descr oder title Attribut zu finden
                nvPicPr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                if nvPicPr is not None:
                    existing_alt = nvPicPr.get('descr') or nvPicPr.get('title')
            
            figure = Figure(
                image_data=image_bytes,
                mime_type=mime_type,
                alt_text=existing_alt if existing_alt and len(existing_alt.strip()) > 3 else None,
                needs_alt_text=not (existing_alt and len(existing_alt.strip()) > 3),
                image_hash=image_hash,
            )
            
            return Block(
                block_type=BlockType.FIGURE,
                reading_order=0,
                figure=figure,
                source_shape_id=str(shape.shape_id),
            )
            
        except Exception as e:
            print(f"⚠️  Bild-Extraktion fehlgeschlagen: {e}")
            return None
    
    def _parse_table(self, shape) -> Optional[Block]:
        """Parst eine Tabelle."""
        try:
            if not hasattr(shape, 'table'):
                return None
            
            pptx_table: PptxTable = shape.table
            
            rows = []
            for row_idx, pptx_row in enumerate(pptx_table.rows):
                cells = []
                for col_idx, pptx_cell in enumerate(pptx_row.cells):
                    # Paragraphen in Zelle
                    paragraphs = []
                    if pptx_cell.text_frame:
                        for pptx_para in pptx_cell.text_frame.paragraphs:
                            para = self._parse_paragraph(pptx_para)
                            if para:
                                paragraphs.append(para)
                    
                    cell = TableCell(
                        paragraphs=paragraphs,
                        is_header=(row_idx == 0),  # Erste Zeile = Header
                    )
                    cells.append(cell)
                
                rows.append(cells)
            
            table = Table(rows=rows)
            
            return Block(
                block_type=BlockType.TABLE,
                reading_order=0,
                table=table,
                source_shape_id=str(shape.shape_id),
            )
            
        except Exception as e:
            print(f"⚠️  Tabellen-Extraktion fehlgeschlagen: {e}")
            return None
    
    def _parse_chart(self, shape) -> Optional[Block]:
        """
        Parst ein Chart/Diagramm.
        
        TODO: Chart-Daten extrahieren für detaillierte Beschreibung
        """
        # Charts werden als Figure behandelt (Screenshot + Alt-Text)
        # Für echte Daten-Extraktion bräuchten wir tieferen XML-Zugriff
        
        figure = Figure(
            image_data=None,  # TODO: Chart als Bild rendern
            mime_type="image/png",
            alt_text=None,
            needs_alt_text=True,
            long_description="[Diagramm - Daten nicht extrahiert]",
        )
        
        return Block(
            block_type=BlockType.FIGURE,
            reading_order=0,
            figure=figure,
            source_shape_id=str(shape.shape_id),
            confidence=0.5,  # Niedrig weil Daten fehlen
        )
    
    def _determine_reading_order(
        self, 
        shapes_with_order: list[tuple]
    ) -> list[Block]:
        """
        Bestimmt die Lesereihenfolge für Blöcke.
        
        Heuristik:
        1. Title-Placeholder zuerst
        2. Dann Top-to-Bottom, Left-to-Right
        3. Bei gleicher Y-Position: Left-to-Right
        """
        if not shapes_with_order:
            return []
        
        def sort_key(item):
            shape, block, bbox = item
            
            # Titel immer zuerst
            if block.block_type == BlockType.HEADING and block.heading_level == 1:
                return (0, 0, 0)
            
            if bbox:
                # Primär: Y-Position (oben nach unten)
                # Sekundär: X-Position (links nach rechts)
                y_bucket = int(bbox.y / 20)  # 20mm Toleranz
                return (1, y_bucket, bbox.x)
            
            return (2, 0, 0)
        
        sorted_items = sorted(shapes_with_order, key=sort_key)
        
        blocks = []
        for order, (shape, block, bbox) in enumerate(sorted_items, 1):
            block.reading_order = order
            blocks.append(block)
        
        return blocks
    
    def _get_bounding_box(self, shape: BaseShape) -> Optional[BoundingBox]:
        """Extrahiert Position und Größe eines Shapes."""
        try:
            return BoundingBox(
                x=self._emu_to_mm(shape.left),
                y=self._emu_to_mm(shape.top),
                width=self._emu_to_mm(shape.width),
                height=self._emu_to_mm(shape.height),
            )
        except:
            return None
    
    def _has_bullet(self, pptx_para) -> bool:
        """Prüft ob ein Absatz Aufzählungszeichen hat."""
        try:
            # python-pptx Bullet-Detection
            pPr = pptx_para._p.pPr
            if pPr is not None:
                buNone = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                if buNone is not None:
                    return False
                buChar = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                buAutoNum = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                return buChar is not None or buAutoNum is not None
        except:
            pass
        return False
    
    @staticmethod
    def _emu_to_mm(emu: int) -> float:
        """Konvertiert EMUs zu Millimetern."""
        return emu / 914400 * 25.4
    
    @staticmethod
    def _emu_to_pt(emu: int) -> float:
        """Konvertiert EMUs zu Points."""
        return emu / 914400 * 72
